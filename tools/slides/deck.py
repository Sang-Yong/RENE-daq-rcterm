# -*- coding: utf-8 -*-
"""
발표자료 공통 뼈대 — 색·글꼴·배치·그림 조각.

2026-08 의 첫 발표자료가 쓰던 디자인을 이어받되, 글머리 기호만 늘어놓는 대신
**그림으로 보여 주는** 조각들을 갖췄다. 슬라이드를 손으로 그리지 않고 코드로
찍는 이유는 이 저장소의 다른 것들과 같다 — 수치가 바뀌면 다시 돌리면 된다.

    from deck import Deck
    d = Deck("제목", "부제")
    d.title(...); d.section(...); d.slide(...)
    d.save("...pptx")
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---- 색 -------------------------------------------------------------
ACCENT  = RGBColor(0x0E, 0x5C, 0x73)   # 청록. 강조와 눈썹줄
ACC2    = RGBColor(0x1B, 0x8A, 0xA8)   # 밝은 청록. 보조
INK     = RGBColor(0x16, 0x1B, 0x22)   # 본문 글자 / 어두운 패널 바탕
INK2    = RGBColor(0x3D, 0x4A, 0x54)
MUTED   = RGBColor(0x8C, 0x98, 0xA4)
FAINT   = RGBColor(0xC3, 0xCE, 0xD6)
PAPER   = RGBColor(0xFF, 0xFF, 0xFF)
PANEL   = RGBColor(0xF2, 0xF5, 0xF7)   # 옅은 카드 바탕
CODEBG  = RGBColor(0x16, 0x1B, 0x22)
CODEFG  = RGBColor(0xD6, 0xDE, 0xE4)
OK      = RGBColor(0x1B, 0x6B, 0x45)
OKBG    = RGBColor(0xE2, 0xF1, 0xE8)
WARN    = RGBColor(0x8A, 0x5A, 0x00)
WARNBG  = RGBColor(0xFA, 0xF0, 0xDA)
CRIT    = RGBColor(0xA3, 0x28, 0x1F)
CRITBG  = RGBColor(0xFA, 0xE8, 0xE6)

MONO = "D2Coding"
SANS = "맑은 고딕"

W, H = 13.333, 7.5
M     = 0.75                    # 좌우 여백
CW    = W - 2 * M               # 본문 폭


def _rgb(v):
    return v if isinstance(v, RGBColor) else RGBColor(*v)


class Deck(object):
    def __init__(self, mono=MONO, sans=SANS):
        self.p = Presentation()
        self.p.slide_width  = Inches(W)
        self.p.slide_height = Inches(H)
        self.blank = self.p.slide_layouts[6]
        self.mono, self.sans = mono, sans

    # ---- 바탕 조각 --------------------------------------------------
    def _s(self):
        return self.p.slides.add_slide(self.blank)

    def box(self, s, x, y, w, h, fill=None, line=None, lw=1.0,
            shape=MSO_SHAPE.RECTANGLE, radius=None):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = _rgb(line); sp.line.width = Pt(lw)
        sp.shadow.inherit = False
        if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
            try:
                sp.adjustments[0] = radius
            except Exception:
                pass
        return sp

    def text(self, s, x, y, w, h, runs, size=15, font=None, color=INK,
             bold=False, align=PP_ALIGN.LEFT, space=6, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
        """runs = 문자열, 또는 (문자열, {폰트 속성}) 목록의 목록(문단별)."""
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        if isinstance(runs, str):
            runs = [runs]
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(space)
            if line_spacing:
                p.line_spacing = line_spacing
            chunks = para if isinstance(para, list) else [(para, {})]
            for txt, opt in chunks:
                r = p.add_run(); r.text = txt
                f = r.font
                f.name = opt.get("font", font or self.sans)
                f.size = Pt(opt.get("size", size))
                f.bold = opt.get("bold", bold)
                f.color.rgb = _rgb(opt.get("color", color))
        return tb

    # ---- 슬라이드 종류 ----------------------------------------------
    def title(self, eyebrow, title, sub, foot):
        s = self._s()
        self.box(s, 0, 0, 0.28, H, fill=ACCENT)
        self.text(s, 1.10, 1.95, 10.6, 0.4, [[(eyebrow, {"font": self.mono, "size": 15,
                                                         "bold": True, "color": ACCENT})]])
        self.text(s, 1.10, 2.42, 10.6, 1.5, [[(title, {"size": 46, "bold": True, "color": INK})]])
        self.box(s, 1.10, 4.20, 1.6, 0.05, fill=ACCENT)
        self.text(s, 1.10, 4.52, 10.6, 1.0,
                  [[(sub, {"size": 19, "color": INK2})]])
        self.text(s, 1.10, 6.05, 10.6, 0.9,
                  [[(l, {"font": self.mono, "size": 13, "color": MUTED})] for l in foot])
        return s

    def section(self, num, title, sub=""):
        s = self._s()
        self.box(s, 0, 0, W, H, fill=INK)
        self.text(s, M + 0.35, 2.35, CW, 1.2,
                  [[(num, {"font": self.mono, "size": 64, "bold": True, "color": ACC2})]])
        self.text(s, M + 0.35, 3.55, CW, 1.0,
                  [[(title, {"size": 38, "bold": True, "color": PAPER})]])
        if sub:
            self.text(s, M + 0.35, 4.60, CW - 1.0, 0.8,
                      [[(sub, {"size": 17, "color": FAINT})]])
        return s

    def head(self, eyebrow, title):
        """본문 슬라이드의 머리. 내용은 y=2.10 부터 쓰면 된다."""
        s = self._s()
        self.text(s, M, 0.62, CW, 0.32,
                  [[(eyebrow, {"font": self.mono, "size": 12, "bold": True, "color": ACCENT})]])
        self.text(s, M, 1.02, CW, 0.85,
                  [[(title, {"size": 30, "bold": True, "color": INK})]])
        self.box(s, M, 1.80, 1.10, 0.045, fill=ACCENT)
        return s

    # ---- 내용 조각 --------------------------------------------------
    def bullets(self, s, x, y, w, items, size=16, gap=9, color=INK, bullet="•  "):
        paras = []
        for it in items:
            if isinstance(it, tuple):
                lead, rest = it
                paras.append([(bullet, {"color": ACCENT, "bold": True, "size": size}),
                              (lead, {"bold": True, "size": size, "color": color}),
                              (rest, {"size": size, "color": color})])
            else:
                paras.append([(bullet, {"color": ACCENT, "bold": True, "size": size}),
                              (it, {"size": size, "color": color})])
        return self.text(s, x, y, w, 0.4, paras, space=gap, line_spacing=1.15)

    def code(self, s, x, y, w, h, lines, size=13.5, pad=0.22):
        self.box(s, x, y, w, h, fill=CODEBG)
        paras = []
        for ln in lines:
            if isinstance(ln, tuple):
                # 실수로 꼬리 쉼표를 붙여 1-튜플이 되는 일이 잦다. 죽지 않게 받는다.
                txt, col = (ln + (CODEFG,))[:2] if len(ln) == 1 else ln
                paras.append([(txt, {"font": self.mono, "size": size, "color": col})])
            else:
                paras.append([(ln, {"font": self.mono, "size": size, "color": CODEFG})])
        self.text(s, x + pad, y + pad * 0.7, w - 2 * pad, h - pad, paras,
                  space=2, line_spacing=1.10)

    def note(self, s, x, y, w, h, label, body, kind="ok"):
        bg, fg = {"ok": (OKBG, OK), "warn": (WARNBG, WARN), "crit": (CRITBG, CRIT),
                  "info": (PANEL, ACCENT)}[kind]
        self.box(s, x, y, w, h, fill=bg)
        self.box(s, x, y, 0.05, h, fill=fg)
        self.text(s, x + 0.28, y + 0.16, w - 0.5, 0.28,
                  [[(label, {"font": self.mono, "size": 11.5, "bold": True, "color": fg})]])
        self.text(s, x + 0.28, y + 0.50, w - 0.5, h - 0.6,
                  [[(body, {"size": 14.5, "color": INK})]], line_spacing=1.15)

    # ---- 그림 조각 --------------------------------------------------
    def chip(self, s, x, y, w, h, text, kind="info", size=12):
        bg, fg = {"ok": (OKBG, OK), "warn": (WARNBG, WARN), "crit": (CRITBG, CRIT),
                  "info": (PANEL, ACCENT), "dark": (INK, PAPER)}[kind]
        self.box(s, x, y, w, h, fill=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.35)
        self.text(s, x, y + (h - 0.22) / 2.0, w, 0.25,
                  [[(text, {"font": self.mono, "size": size, "bold": True, "color": fg})]],
                  align=PP_ALIGN.CENTER)

    def node(self, s, x, y, w, h, title, sub="", kind="info", tsize=14, ssize=11):
        """흐름도의 상자 하나."""
        bg, fg, ln = {"ok": (OKBG, OK, OK), "warn": (WARNBG, WARN, WARN),
                      "crit": (CRITBG, CRIT, CRIT), "info": (PAPER, INK, FAINT),
                      "dark": (INK, PAPER, INK)}[kind]
        self.box(s, x, y, w, h, fill=bg, line=ln, lw=1.25)
        ty = y + (h - (0.30 if sub else 0.22)) / 2.0 - (0.10 if sub else 0)
        self.text(s, x + 0.10, ty, w - 0.20, 0.30,
                  [[(title, {"font": self.mono, "size": tsize, "bold": True, "color": fg})]],
                  align=PP_ALIGN.CENTER)
        if sub:
            self.text(s, x + 0.10, ty + 0.30, w - 0.20, 0.28,
                      [[(sub, {"size": ssize, "color": MUTED if kind != "dark" else FAINT})]],
                      align=PP_ALIGN.CENTER)

    def arrow(self, s, x, y, w, label="", color=ACCENT, size=10.5, down=False):
        if down:
            self.box(s, x, y, 0.035, w, fill=color)
            tip = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                     Inches(x - 0.075), Inches(y + w), Inches(0.185), Inches(0.15))
            tip.rotation = 180
        else:
            self.box(s, x, y, w, 0.035, fill=color)
            tip = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                     Inches(x + w), Inches(y - 0.075), Inches(0.15), Inches(0.185))
            tip.rotation = 90
        tip.fill.solid(); tip.fill.fore_color.rgb = _rgb(color)
        tip.line.fill.background(); tip.shadow.inherit = False
        if label:
            if down:
                self.text(s, x + 0.16, y + w / 2 - 0.12, 2.0, 0.25,
                          [[(label, {"font": self.mono, "size": size, "color": color})]])
            else:
                self.text(s, x, y - 0.34, w, 0.25,
                          [[(label, {"font": self.mono, "size": size, "color": color})]],
                          align=PP_ALIGN.CENTER)

    def flow(self, s, y, nodes, h=0.95, gap=0.42, x0=None, wtot=None, labels=None):
        """가로 흐름도. nodes = [(제목, 부제, 종류), ...]"""
        n = len(nodes)
        wtot = wtot if wtot is not None else CW
        x0 = x0 if x0 is not None else M
        bw = (wtot - gap * (n - 1)) / n
        xs = []
        for i, nd in enumerate(nodes):
            x = x0 + i * (bw + gap)
            xs.append(x)
            t, sub, kind = (nd + ("info",))[:3] if len(nd) < 3 else nd
            self.node(s, x, y, bw, h, t, sub, kind)
            if i < n - 1:
                lab = labels[i] if labels and i < len(labels) else ""
                self.arrow(s, x + bw + 0.06, y + h / 2 - 0.02, gap - 0.24, lab)
        return xs, bw

    def bar(self, s, x, y, w, h, frac, color=ACCENT, bg=PANEL, label="", lsize=12):
        self.box(s, x, y, w, h, fill=bg)
        if frac > 0:
            self.box(s, x, y, max(0.03, w * min(1.0, frac)), h, fill=color)
        if label:
            self.text(s, x + w + 0.14, y + (h - 0.2) / 2, 3.2, 0.25,
                      [[(label, {"font": self.mono, "size": lsize, "color": INK})]])

    def table(self, s, x, y, w, cols, rows, widths=None, rh=0.36, size=13,
              head_bg=INK, head_fg=PAPER, zebra=True):
        n = len(cols)
        widths = widths or [w / n] * n
        cx = x
        for j, c in enumerate(cols):
            self.box(s, cx, y, widths[j], rh, fill=head_bg)
            self.text(s, cx + 0.12, y + (rh - 0.20) / 2, widths[j] - 0.2, 0.24,
                      [[(c, {"font": self.mono, "size": size - 1.5, "bold": True, "color": head_fg})]])
            cx += widths[j]
        for i, row in enumerate(rows):
            ry = y + rh * (i + 1)
            cx = x
            if zebra and i % 2 == 1:
                self.box(s, x, ry, w, rh, fill=PANEL)
            for j, cell in enumerate(row):
                txt, opt = (cell, {}) if isinstance(cell, str) else cell
                self.text(s, cx + 0.12, ry + (rh - 0.20) / 2, widths[j] - 0.2, 0.24,
                          [[(txt, {"size": opt.get("size", size),
                                   "font": opt.get("font", self.sans),
                                   "bold": opt.get("bold", False),
                                   "color": opt.get("color", INK)})]])
                cx += widths[j]
        return y + rh * (len(rows) + 1)

    def foot(self, s, text):
        self.text(s, M, H - 0.62, CW, 0.3,
                  [[(text, {"font": self.mono, "size": 10.5, "color": MUTED})]])

    def save(self, path):
        self.p.save(path)
        return path
