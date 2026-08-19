# -*- coding: utf-8 -*-
"""배치 점검 — 렌더러 없이 슬라이드가 깨졌는지 본다.

LibreOffice 가 없는 서버에서 만들다 보니 눈으로 확인할 수가 없다. 대신
① 슬라이드 밖으로 나간 도형 ② 글이 상자보다 길어 넘칠 것 같은 곳
③ 글상자끼리 겹치는 곳을 수치로 잡는다. 추정이라 완벽하지 않지만,
제목이 잘리거나 표가 밖으로 나가는 종류의 사고는 거의 다 걸린다.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0


def width_em(ch):
    o = ord(ch)
    # 한글·한자·전각은 1 em, 그 밖은 대략 0.5 em (D2Coding 은 0.6)
    if 0x1100 <= o <= 0x11FF or 0x3000 <= o <= 0x303F or \
       0xAC00 <= o <= 0xD7A3 or 0x4E00 <= o <= 0x9FFF or 0xFF00 <= o <= 0xFF60:
        return 1.0
    return 0.52


def para_lines(text, size_pt, box_w_in):
    """글 한 문단이 몇 줄이 될지 추정."""
    if not text:
        return 1
    cap_em = box_w_in * 72.0 / max(size_pt, 1.0)
    lines, cur = 1, 0.0
    for ch in text:
        if ch == "\n":
            lines += 1; cur = 0.0; continue
        cur += width_em(ch)
        if cur > cap_em:
            lines += 1; cur = width_em(ch)
    return lines


def audit(path, verbose=False):
    p = Presentation(path)
    SW, SH = p.slide_width / EMU, p.slide_height / EMU
    problems = []
    for i, s in enumerate(p.slides, 1):
        boxes = []
        for sh in s.shapes:
            L, T = sh.left / EMU, sh.top / EMU
            R, B = L + sh.width / EMU, T + sh.height / EMU
            if L < -0.02 or T < -0.02 or R > SW + 0.02 or B > SH + 0.02:
                problems.append((i, "슬라이드 밖", "%s  L%.2f T%.2f R%.2f B%.2f"
                                 % (sh.name, L, T, R, B)))
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            txt = tf.text
            if not txt.strip():
                continue
            need = 0.0
            for para in tf.paragraphs:
                sz = None
                for r in para.runs:
                    if r.font.size:
                        sz = max(sz or 0, r.font.size.pt)
                sz = sz or 18.0
                ls = para.line_spacing or 1.0
                if not isinstance(ls, float):
                    ls = 1.0
                n = para_lines("".join(r.text for r in para.runs), sz, sh.width / EMU)
                need += n * sz * 1.22 * ls / 72.0
                need += (para.space_after.pt / 72.0) if para.space_after else 0.0
            have = sh.height / EMU
            # 글상자는 자동으로 늘어나므로, 아래 여백까지 먹으면 문제로 본다
            if T + need > SH - 0.10:
                problems.append((i, "아래로 넘침",
                                 "%s  필요 %.2f\" 시작 %.2f\" -> 끝 %.2f\" (슬라이드 %.2f\")"
                                 % (sh.name, need, T, T + need, SH)))
            boxes.append((sh.name, L, T, L + sh.width / EMU, T + max(need, have), txt[:24]))
        for a in range(len(boxes)):
            for b in range(a + 1, len(boxes)):
                n1, l1, t1, r1, b1, x1 = boxes[a]
                n2, l2, t2, r2, b2, x2 = boxes[b]
                ov = min(r1, r2) - max(l1, l2), min(b1, b2) - max(t1, t2)
                if ov[0] > 0.25 and ov[1] > 0.18:
                    problems.append((i, "글끼리 겹침",
                                     "'%s' x '%s'  겹침 %.2f x %.2f\"" % (x1, x2, ov[0], ov[1])))
    return SW, SH, len(p.slides._sldIdLst), problems


if __name__ == "__main__":
    for f in sys.argv[1:]:
        SW, SH, n, probs = audit(f)
        print("\n=== %s ===  %d장  %.2f x %.2f in" % (f.split("/")[-1], n, SW, SH))
        if not probs:
            print("  문제 없음")
        else:
            cur = None
            for sl, kind, det in probs:
                if sl != cur:
                    print("  --- %d장 ---" % sl); cur = sl
                print("    [%s] %s" % (kind, det))
            print("  총 %d건" % len(probs))
