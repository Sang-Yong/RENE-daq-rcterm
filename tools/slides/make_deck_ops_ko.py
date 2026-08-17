#!/usr/bin/env python3
"""RENE DAQ run control - presentation generator."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette -------------------------------------------------------------
INK      = RGBColor(0x16, 0x1B, 0x22)
MUTED    = RGBColor(0x5A, 0x66, 0x72)
FAINT    = RGBColor(0x8C, 0x98, 0xA4)
ACCENT   = RGBColor(0x0E, 0x5C, 0x73)   # deep teal
ACCENT2  = RGBColor(0xC2, 0x5E, 0x00)   # burnt orange
GOOD     = RGBColor(0x1B, 0x6B, 0x45)
BAD      = RGBColor(0xA3, 0x28, 0x20)
RULE     = RGBColor(0xD8, 0xDF, 0xE4)
PANEL    = RGBColor(0xF2, 0xF5, 0xF7)
CODEBG   = RGBColor(0x16, 0x1B, 0x22)
CODEFG   = RGBColor(0xD6, 0xDE, 0xE4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "맑은 고딕"
MONO = "D2Coding"          # 없으면 PowerPoint 가 대체한다
SANS_EA = "맑은 고딕"
MONO_EA = "맑은 고딕"


def _set_font(run, name):
    """latin 뿐 아니라 ea(한중일) / cs 타이프페이스까지 지정한다.
    이걸 안 하면 한글이 PowerPoint 기본 글꼴로 제멋대로 바뀐다."""
    from pptx.oxml.ns import qn
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = name if name != MONO else MONO_EA
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", ea)

W, H = Inches(13.333), Inches(7.5)
M    = Inches(0.75)                      # margin

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text, size, color, bold=False, font=SANS, space_before=0,
       space_after=4, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold = Pt(size), bold
    _set_font(r, font)
    r.font.color.rgb = color
    return p


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def header(slide, title, kicker=None):
    y = M
    if kicker:
        tf = _tb(slide, M, y - Inches(0.06), W - 2 * M, Inches(0.3))
        _p(tf, kicker.upper(), 12, ACCENT, bold=True, font=MONO, first=True)
        y += Inches(0.34)
    tf = _tb(slide, M, y, W - 2 * M, Inches(0.8))
    _p(tf, title, 30, INK, bold=True, first=True)
    rect(slide, M, y + Inches(0.72), Inches(1.1), Pt(3), ACCENT)
    return y + Inches(1.05)


def bullets(slide, items, top, size=17, width=None, left=None, gap=9):
    """items: list of (text, kind) where kind in {'', 'sub', 'good', 'bad', 'code'}"""
    left = left or M
    width = width or (W - 2 * M)
    tf = _tb(slide, left, top, width, H - top - M)
    # 대략적인 줄 수로 아래쪽 y 를 되돌려 준다. 다음 요소를 이어 붙이기 위한 것.
    chars_per_line = max(20, int(width / Inches(1.0) * 11.6 * (17.0 / size)))
    used = Emu(0)
    first = True

    def _w(t):
        # 한글은 라틴 문자보다 넓다. 줄 수 추정에서 1.8 배로 센다.
        return sum(1.8 if '\uac00' <= c <= '\ud7a3' or '\u3130' <= c <= '\u318f'
                   else 1.0 for c in t)

    for text, kind in items:
        nlines = max(1, -(-int(_w(text)) // chars_per_line))
        fs = size - 3 if kind in ("sub", "code") else size
        used += Emu(int(Pt(fs * 1.28) * nlines + Pt(gap)))
    for text, kind in items:
        if kind == "sub":
            _p(tf, "     " + text, size - 3, MUTED, space_after=gap - 2, first=first)
        elif kind == "good":
            _p(tf, "✓  " + text, size, GOOD, bold=True, space_after=gap, first=first)
        elif kind == "bad":
            _p(tf, "✗  " + text, size, BAD, bold=True, space_after=gap, first=first)
        elif kind == "code":
            _p(tf, text, size - 3, ACCENT, font=MONO, space_after=gap, first=first)
        elif kind == "head":
            _p(tf, text, size, ACCENT, bold=True, space_before=8, space_after=gap, first=first)
        else:
            _p(tf, "•  " + text, size, INK, space_after=gap, first=first)
        first = False
    return top + used


def codebox(slide, lines, x, y, w, h=None, size=13, title=None):
    lh = Pt(size * 1.45)
    h = h or Emu(int(lh * (len(lines) + 1.4)))
    rect(slide, x, y, w, h, CODEBG)
    tf = _tb(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), h - Inches(0.3))
    first = True
    if title:
        _p(tf, title, size - 1, FAINT, font=MONO, space_after=5, first=True)
        first = False
    for ln in lines:
        col = CODEFG
        if ln.startswith("#") or ln.startswith("//"):
            col = FAINT
        _p(tf, ln, size, col, font=MONO, space_after=1, first=first)
        first = False
    return y + h


def metric_row(slide, items, y, h=Inches(1.5)):
    """items: list of (big, label, color)"""
    n = len(items)
    gap = Inches(0.25)
    w = (W - 2 * M - gap * (n - 1)) / n
    for i, (big, label, col) in enumerate(items):
        x = M + i * (w + gap)
        rect(slide, x, y, w, h, PANEL)
        rect(slide, x, y, Pt(4), h, col)
        tf = _tb(slide, x + Inches(0.28), y + Inches(0.22), w - Inches(0.5), h - Inches(0.4))
        _p(tf, big, 34, col, bold=True, space_after=2, first=True)
        _p(tf, label, 12.5, MUTED)
    return y + h


def table(slide, rows, top, widths, size=14, head=True):
    x = M
    y = top
    rowh = Inches(0.42)
    for ri, row in enumerate(rows):
        cx = x
        if ri == 0 and head:
            rect(slide, M, y, sum(widths), rowh, PANEL)
        for ci, cell in enumerate(row):
            tf = _tb(slide, cx + Inches(0.14), y + Inches(0.08), widths[ci] - Inches(0.2), rowh)
            col = MUTED if (ri == 0 and head) else INK
            bold = (ri == 0 and head)
            fnt = SANS
            if cell.startswith("`") and cell.endswith("`"):
                cell = cell[1:-1]; fnt = MONO
            if cell.startswith("+"):
                col = GOOD; cell = cell[1:]
            elif cell.startswith("-"):
                col = BAD; cell = cell[1:]
            _p(tf, cell, size, col, bold=bold, font=fnt, first=True)
            cx += widths[ci]
        y += rowh
        if ri == 0 and head:
            rect(slide, M, y, sum(widths), Pt(1), RULE)
        elif ri < len(rows) - 1:
            rect(slide, M, y, sum(widths), Pt(0.75), RULE)
    return y


def footer(slide, text):
    tf = _tb(slide, M, H - Inches(0.62), W - 2 * M, Inches(0.3))
    _p(tf, text, 11.5, FAINT, first=True)


def new(kicker=None, title=None):
    s = prs.slides.add_slide(BLANK)
    top = header(s, title, kicker) if title else M
    return s, top


def divider(number, title, sub):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, ACCENT)
    tf = _tb(s, M, H / 2 - Inches(1.0), W - 2 * M, Inches(2.2))
    _p(tf, f"{number:02d}", 15, RGBColor(0x8F, 0xC5, 0xD4), bold=True, font=MONO, first=True)
    _p(tf, title, 40, WHITE, bold=True, space_before=4, space_after=10)
    _p(tf, sub, 17, RGBColor(0xC6, 0xDF, 0xE8))
    return s


# =========================================================================
#  운영자용 슬라이드 본문
# =========================================================================

# ---- 1. 표지 ----
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, Inches(0.28), H, ACCENT)
tf = _tb(s, Inches(1.1), Inches(2.0), Inches(10.6), Inches(3.2))
_p(tf, "RENE / CUPDAQ", 15, ACCENT, bold=True, font=MONO, first=True)
_p(tf, "DAQ 운용 안내", 46, INK, bold=True, space_before=6, space_after=16)
_p(tf, "화면 보는 법 · 시작과 종료 · 이상할 때 대처", 20, MUTED)
tf2 = _tb(s, Inches(1.1), H - Inches(1.5), Inches(10.6), Inches(0.9))
_p(tf2, "rcterm  ·  rcsupervisor  ·  postrun", 13, FAINT, font=MONO, first=True)
_p(tf2, "2026년 8월   ·   영광 사이트", 13, FAINT, font=MONO)

# ---- 2. 한눈에 ----
s, top = new("한눈에", "무엇이 어디서 도는가")
codebox(s, [
    "rcsupervisor      24시간마다 런을 교체하고, 10분마다 이상을 점검한다",
    "     │                이것 하나만 띄우면 나머지는 알아서 뜬다",
    "     ▼",
    "  rcterm          런 1개를 진행한다  boot → config → start → 감시 → end",
    "     │",
    "     ▼  executedaq.sh 가 띄운다",
    "  daq(FADC)   daq(SADC)   tcb        실제로 데이터를 받는 프로그램",
    "",
    "  postrun         수집을 3분 뒤에서 따라가며 병합·분석파일을 만든다",
], M, top, W - 2 * M, size=14.5)
bullets(s, [
    ("운영자가 직접 띄우는 것은 rcsupervisor 하나뿐이다. daq / tcb 를 손으로 "
     "실행할 일은 없다.", ""),
    ("런은 24시간마다 자동으로 교체된다. 교체할 때 10~40초 정도 수집이 끊기는데, "
     "프로토콜에 '수집 중 run 번호 변경' 명령이 없어 어쩔 수 없다.", ""),
    ("run 4288→4289→4290 으로 두 번 교체되는 동안 카탈로그 기록·이벤트 수가 "
     "모두 정상이었고, 서브런은 FADC·SADC·Merged·PRD 가 1440 으로 일치했다.", ""),
], Inches(4.95), size=15.5)

# ---- 3. 화면 보는 법 ----
s, top = new("화면", "pane 다섯이 각각 무엇인가")
codebox(s, [
    "+-------------------------+-------------------+",
    "| DAQ Run Status(monitor) |                   |  monitor    지금 수집 상태",
    "+-------------------------+   work space      |  supervisor 감시자 로그",
    "| supervisor              |                   |  postrun    병합·분석 진행",
    "+-------------------------+-------------------+  work space 자유롭게 작업",
    "| postrun                 | dataflow:         |  dataflow   데이터가 어디까지",
    "|                         |  /Data_ssd(RAW)   |             옮겨졌나",
    "|                         |  ->/data(PRD)     |",
    "|                         |  ->khu(backup)    |",
    "|                         |  ->scratch(save)  |",
    "+-------------------------+-------------------+",
], M, top, W - 2 * M, size=14)
bullets(s, [
    ("work space 를 뺀 넷은 보기만 하는 화면이다. 명령은 work space 에서 친다.", ""),
    ("monitor 는 heartbeat 파일을 읽어 그리는 것이라 Ctrl-C 로 꺼도 DAQ 는 "
     "계속 돈다. 몇 개를 띄워도 무방하다.", ""),
    ("dataflow 의 제목이 곧 데이터가 지나는 길이다 — 수집한 자리에서 "
     "백업 대기, 경희대 백업, 마지막으로 장기보관.", ""),
    ("비율이 틀어지면 Ctrl-B 다음 = 를 누르면 원래대로 돌아온다.", ""),
], Inches(4.9), size=16)

# ---- 4. monitor 읽는 법 ----
s, top = new("화면", "monitor — 이 화면만 보면 된다")
_c4 = [
    "======================================================================",
    "  RENE / CUPDAQ   Run Monitor   (heartbeat viewer, read-only)",
    "======================================================================",
    "        Current Time : 2026-08-17 05:35:31",
    "          Run Number : 4290 / 42          ← run 번호 / 서브런 번호",
    "           DAQ State : Running            ← 이게 Running 이어야 정상",
    "            DAQ Time : 00:42:52           ← 이 런이 수집한 시간",
    "        Total Events : 5223856",
    "  ------------------------------------------------------------------",
    "        DAQ          Events       Rate[Hz]     Average[Hz]",
    "  ------------------------------------------------------------------",
    "       FADCDAQ        2612144         997.3          1015.4",
    "       SADCDAQ        2611712        1011.3          1015.3",
    "  ------------------------------------------------------------------",
    "  hb 0s                                  ← 가장 중요한 한 줄",
]
_y4 = codebox(s, _c4, M, top, W - 2 * M, size=12)
bullets(s, [
    ("hb 는 heartbeat 나이다. 0~2초면 정상이고, 이 숫자가 계속 커지면 rcterm 이 "
     "멈춘 것이다. 화면이 얼어붙은 것과 DAQ 가 멈춘 것을 이걸로 구분한다.", ""),
    ("두 ADC 의 Rate 가 비슷하고 Average 근처에 있으면 정상이다.", ""),
], _y4 + Inches(0.12), size=15)

# ---- 5. 시작하기 ----
s, top = new("조작", "시작하기")
_y5 = codebox(s, [
    "cd /home/frontend/DAQ/RENE-daq-rcterm",
    "",
    "scripts/daq-tmux.sh                # 화면만 만든다. 하드웨어는 건드리지 않는다",
    "scripts/daq-tmux.sh --start        # DAQ 가 꺼져 있으면 같이 기동한다",
], M, top, W - 2 * M, size=15)
y = bullets(s, [
    ("--start 없이 실행하면 감시자 명령을 pane 에 '입력만' 해 둔다. 확인하고 "
     "Enter 를 치면 된다. 실수로 이중 기동하는 것을 막기 위해서다.", ""),
    ("이미 DAQ 가 돌고 있으면 절대 다시 띄우지 않는다. 그 경우 감시자 pane 은 "
     "로그만 따라간다.", ""),
    ("이미 daq 세션이 있으면 화면을 다시 짜지 않고 그냥 붙는다. 살아있는 런을 "
     "실수로 흔들지 않기 위해서다.", ""),
], _y5 + Inches(0.2), size=16)
codebox(s, [
    "# 지금 DAQ 가 도는지 확인",
    "pgrep -x rcsupervisor        # 숫자가 나오면 돌고 있다",
    "cat /Data/LOG/rcterm.hb      # time= 이 계속 바뀌면 살아있다",
], M, y + Inches(0.1), W - 2 * M, size=14)

# ---- 6. 화면 조작 ----
s, top = new("조작", "붙었다 떼기, 화면 옮겨다니기")
y = table(s, [
    ["하려는 것", "누르는 키 / 명령"],
    ["화면에 붙기", "`tmux attach -t daq`"],
    ["화면에서 떼기 (DAQ 는 계속 돈다)", "`Ctrl-B  다음  D`"],
    ["pane 사이 이동", "`Ctrl-B  다음  방향키`"],
    ["pane 번호 크게 보기", "`Ctrl-B  다음  q`"],
    ["pane 비율 원래대로", "`Ctrl-B  다음  =`"],
    ["pane 경계 옮기기", "`Ctrl-B  다음  H J K L` (연타 가능)"],
    ["한 pane 만 크게 (다시 누르면 원래대로)", "`Ctrl-B  다음  z`"],
], top, [Inches(6.3), Inches(5.5)])
bullets(s, [
    ("떼어놓아도(detach) DAQ 는 계속 돈다. 터미널 창을 닫아도, ssh 가 끊겨도, "
     "화면이 죽어도 마찬가지다. 그것이 tmux 를 쓰는 이유다.", ""),
], y + Inches(0.35), size=16)

# ---- 7. 안전하게 멈추기 ----
s, top = new("조작 · 중요", "안전하게 멈추기")
_y7 = codebox(s, [
    "kill -TERM $(pgrep -x rcsupervisor)",
    "",
    "# 감시자 → rcterm → ENDRUN → EXIT 순서로 얌전히 내려간다",
    "# DB 에 런이 기록되고 끝난다",
], M, top, W - 2 * M, size=15)
y = bullets(s, [
    ("신호는 감시자 PID 에만 보낸다. 이것만 지키면 된다.", "good"),
], _y7 + Inches(0.2), size=17)
_y7b = bullets(s, [
    ("프로세스 그룹에 보내지 말 것 (kill -TERM -PID)", "bad"),
    ("daq / tcb / merger 를 직접 죽이지 말 것", "bad"),
    ("rcterm 에 직접 보내지 말 것 — 감시자가 장애로 보고 새 런을 띄운다", "bad"),
], y + Inches(0.1), size=17)
bullets(s, [
    ("daq / tcb / merger 는 rcterm 의 자식이 아니다. 그룹에 신호를 보내면 "
     "쓰기 도중에 끊겨 마지막 파일이 상한다. 정말 강제로 정리해야 하면 "
     "scripts/killdaq.sh 를 쓴다.", ""),
], _y7b + Inches(0.15), size=15.5)

# ---- 8. 후처리 확인 ----
s, top = new("후처리", "병합·분석이 따라오고 있는가")
_y8 = codebox(s, [
    "[2026-08-17 05:34:20] Merging FADC Subrun 38 ...",
    "[2026-08-17 05:34:48]  -> Producing Subrun 38 (병렬 슬롯 최대 3)",
    "[2026-08-17 05:34:49] Merging FADC Subrun 39 ...",
    "[2026-08-17 05:34:57]  [OK] Producing Done : Subrun 38 (총 처리시간: 9초)",
], M, top, W - 2 * M, size=14, title="postrun pane")
y = bullets(s, [
    ("서브런 하나가 60초 분량인데 처리에 약 28초 걸린다. 즉 따라잡는다.", ""),
    ("monitor 의 서브런 번호와 postrun 의 서브런 번호 차이가 계속 벌어지면 "
     "뒤처지고 있는 것이다.", ""),
], _y8 + Inches(0.2), size=16)
codebox(s, [
    "# 어디까지 처리됐는지",
    "ls /Data_ssd/RAW/004290/PRD/*.root | wc -l",
    "",
    "# 런 전체를 점검 (기존 도구 그대로 쓸 수 있다)",
    "cd /home/frontend/DAQ/DAQ_cup/production/Shell && ./audit_run.sh 004289",
], M, y + Inches(0.1), W - 2 * M, size=14)

# ---- 9. 이상할 때 ----
s, top = new("대처", "이상하면 이 순서로")
y = table(s, [
    ["증상", "먼저 볼 것", "대처"],
    ["monitor 의 hb 가 계속 커진다", "`rcterm 이 멈춤`", "감시자가 5분 뒤 자동 재시작"],
    ["화면이 안 바뀐다", "`hb 는 0s`", "화면만 문제. DAQ 는 정상"],
    ["Rate 가 0 근처", "`TCB 로그`", "트리거 확인. 배경런이면 정상일 수 있음"],
    ["FADC 와 SADC 개수가 다르다", "`postrun 경고`", "한쪽 누락. 로그 확인"],
    ["`a DAQ is already listening`", "`남은 프로세스`", "`scripts/killdaq.sh` 후 재시작"],
    ["postrun 이 계속 뒤처진다", "`디스크 여유`", "`--jobs` 는 올려도 소용없다"],
], top, [Inches(4.6), Inches(3.2), Inches(4.0)])
bullets(s, [
    ("대부분은 감시자가 스스로 복구한다. 10분마다 점검해서 이상하면 런을 "
     "정리하고 새 번호로 다시 띄운다. 연속 5회 실패하면 감시자도 멈추므로, "
     "그때는 사람이 봐야 한다.", ""),
], y + Inches(0.35), size=16)

# ---- 10. 자주 보는 메시지 ----
s, top = new("대처", "자주 보는 메시지와 뜻")
y = table(s, [
    ["메시지", "뜻"],
    ["`[SUP] health OK  run=... state=Running`", "+10분 점검 통과. 정상"],
    ["`[SUP] cycle N finished : exit=code 0`", "+런이 정상으로 끝났다"],
    ["`[SUP] recovering : cleaning up`", "이상을 감지해 정리 중"],
    ["`[Pre-Check] 0-byte 파일 없음. 양호.`", "+후처리 입력 파일 정상"],
    ["`[CORRUPTION DETECTED] ZOMBIE FILE`", "-그 서브런을 읽을 수 없어 건너뜀"],
    ["`[FATAL] a DAQ is already listening`", "-이전 DAQ 가 남아있다"],
    ["`[FATAL] cannot connect to TCB`", "-부팅 실패. TCB 로그를 볼 것"],
], top, [Inches(6.6), Inches(5.2)])
bullets(s, [
    ("[FATAL] 로 끝나면 종료코드가 갈린다. 1 은 설정 오류라 다시 띄워도 소용없고, "
     "2 는 런 실패라 감시자가 정리 후 재시도한다.", ""),
], y + Inches(0.3), size=16)

# ---- 11. 치트시트 ----
s, top = new("요약", "이 한 장만 기억하면 된다")
codebox(s, [
    "시작        scripts/daq-tmux.sh --start",
    "붙기        tmux attach -t daq",
    "떼기        Ctrl-B  다음  D              ← DAQ 는 계속 돈다",
    "비율복원    Ctrl-B  다음  =",
    "",
    "멈추기      kill -TERM $(pgrep -x rcsupervisor)      ← 감시자에만",
    "강제정리    scripts/killdaq.sh                       ← 정말 필요할 때만",
    "",
    "살아있나    pgrep -x rcsupervisor",
    "상태보기    cat /Data/LOG/rcterm.hb        ← time= 이 바뀌면 정상",
    "감시자로그  tail -f /Data/LOG/rcsupervisor.log",
    "후처리점검  cd .../production/Shell && ./audit_run.sh 004289",
], M, top, W - 2 * M, size=14.5)
bullets(s, [
    ("가장 중요한 두 가지 — monitor 의 hb 가 0s 근처인지 보고, 멈출 때는 "
     "감시자 PID 에만 신호를 보낸다.", ""),
], Inches(6.15), size=16)

# ---- 12. 마무리 ----
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, PANEL)
rect(s, 0, 0, Inches(0.28), H, ACCENT)
tf = _tb(s, Inches(1.2), Inches(1.7), Inches(10.8), Inches(4.2))
_p(tf, "정리", 13, ACCENT, bold=True, font=MONO, first=True)
_p(tf, "띄우는 것은 하나, 보는 것은 hb,", 32, INK, bold=True,
   space_before=14, space_after=2)
_p(tf, "멈출 때는 감시자에만.", 32, INK, bold=True, space_after=24)
_p(tf, "자세한 내용", 14, ACCENT, bold=True)
_p(tf, "docs/MANUAL.md        운용 상세", 15, MUTED, font=MONO, space_after=2)
_p(tf, "docs/POSTRUN.md       병합·분석 파이프라인", 15, MUTED, font=MONO, space_after=2)
_p(tf, "CLAUDE.md             설계 근거와 이력", 15, MUTED, font=MONO)

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "ops.pptx"
prs.save(out)
print(f"{len(prs.slides._sldIdLst)} slides -> {out}")
