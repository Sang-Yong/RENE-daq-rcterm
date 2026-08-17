#!/usr/bin/env python3
"""RENE DAQ run control - presentation generator."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette -------------------------------------------------------------
INK      = RGBColor(0x16, 0x1B, 0x22)
MUTED    = RGBColor(0x5A, 0x66, 0x72)
FAINT    = RGBColor(0x8C, 0x98, 0xA4)
ACCENT   = RGBColor(0x0E, 0x5C, 0x73)   # deep teal
ACCENT2  = RGBColor(0xC2, 0x5E, 0x00)   # burnt orange
GOOD     = RGBColor(0x1B, 0x6B, 0x45)
BAD      = RGBColor(0xA3, 0x28, 0x20)
RULE     = RGBColor(0xD8, 0xDF, 0xE4)
PANEL    = RGBColor(0xF2, 0xF5, 0xF7)
CODEBG   = RGBColor(0x16, 0x1B, 0x22)
CODEFG   = RGBColor(0xD6, 0xDE, 0xE4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M    = Inches(0.75)                      # margin

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text, size, color, bold=False, font=SANS, space_before=0,
       space_after=4, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(size), bold, font
    r.font.color.rgb = color
    return p


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def header(slide, title, kicker=None):
    y = M
    if kicker:
        tf = _tb(slide, M, y - Inches(0.06), W - 2 * M, Inches(0.3))
        _p(tf, kicker.upper(), 12, ACCENT, bold=True, font=MONO, first=True)
        y += Inches(0.34)
    tf = _tb(slide, M, y, W - 2 * M, Inches(0.8))
    _p(tf, title, 30, INK, bold=True, first=True)
    rect(slide, M, y + Inches(0.72), Inches(1.1), Pt(3), ACCENT)
    return y + Inches(1.05)


def bullets(slide, items, top, size=17, width=None, left=None, gap=9):
    """items: list of (text, kind) where kind in {'', 'sub', 'good', 'bad', 'code'}"""
    left = left or M
    width = width or (W - 2 * M)
    tf = _tb(slide, left, top, width, H - top - M)
    # 대략적인 줄 수로 아래쪽 y 를 되돌려 준다. 다음 요소를 이어 붙이기 위한 것.
    chars_per_line = max(20, int(width / Inches(1.0) * 11.6 * (17.0 / size)))
    used = Emu(0)
    first = True
    for text, kind in items:
        nlines = max(1, -(-len(text) // chars_per_line))
        fs = size - 3 if kind in ("sub", "code") else size
        used += Emu(int(Pt(fs * 1.28) * nlines + Pt(gap)))
    for text, kind in items:
        if kind == "sub":
            _p(tf, "     " + text, size - 3, MUTED, space_after=gap - 2, first=first)
        elif kind == "good":
            _p(tf, "✓  " + text, size, GOOD, bold=True, space_after=gap, first=first)
        elif kind == "bad":
            _p(tf, "✗  " + text, size, BAD, bold=True, space_after=gap, first=first)
        elif kind == "code":
            _p(tf, text, size - 3, ACCENT, font=MONO, space_after=gap, first=first)
        elif kind == "head":
            _p(tf, text, size, ACCENT, bold=True, space_before=8, space_after=gap, first=first)
        else:
            _p(tf, "•  " + text, size, INK, space_after=gap, first=first)
        first = False
    return top + used


def codebox(slide, lines, x, y, w, h=None, size=13, title=None):
    lh = Pt(size * 1.45)
    h = h or Emu(int(lh * (len(lines) + 1.4)))
    rect(slide, x, y, w, h, CODEBG)
    tf = _tb(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), h - Inches(0.3))
    first = True
    if title:
        _p(tf, title, size - 1, FAINT, font=MONO, space_after=5, first=True)
        first = False
    for ln in lines:
        col = CODEFG
        if ln.startswith("#") or ln.startswith("//"):
            col = FAINT
        _p(tf, ln, size, col, font=MONO, space_after=1, first=first)
        first = False
    return y + h


def metric_row(slide, items, y, h=Inches(1.5)):
    """items: list of (big, label, color)"""
    n = len(items)
    gap = Inches(0.25)
    w = (W - 2 * M - gap * (n - 1)) / n
    for i, (big, label, col) in enumerate(items):
        x = M + i * (w + gap)
        rect(slide, x, y, w, h, PANEL)
        rect(slide, x, y, Pt(4), h, col)
        tf = _tb(slide, x + Inches(0.28), y + Inches(0.22), w - Inches(0.5), h - Inches(0.4))
        _p(tf, big, 34, col, bold=True, space_after=2, first=True)
        _p(tf, label, 12.5, MUTED)
    return y + h


def table(slide, rows, top, widths, size=14, head=True):
    x = M
    y = top
    rowh = Inches(0.42)
    for ri, row in enumerate(rows):
        cx = x
        if ri == 0 and head:
            rect(slide, M, y, sum(widths), rowh, PANEL)
        for ci, cell in enumerate(row):
            tf = _tb(slide, cx + Inches(0.14), y + Inches(0.08), widths[ci] - Inches(0.2), rowh)
            col = MUTED if (ri == 0 and head) else INK
            bold = (ri == 0 and head)
            fnt = SANS
            if cell.startswith("`") and cell.endswith("`"):
                cell = cell[1:-1]; fnt = MONO
            if cell.startswith("+"):
                col = GOOD; cell = cell[1:]
            elif cell.startswith("-"):
                col = BAD; cell = cell[1:]
            _p(tf, cell, size, col, bold=bold, font=fnt, first=True)
            cx += widths[ci]
        y += rowh
        if ri == 0 and head:
            rect(slide, M, y, sum(widths), Pt(1), RULE)
        elif ri < len(rows) - 1:
            rect(slide, M, y, sum(widths), Pt(0.75), RULE)
    return y


def footer(slide, text):
    tf = _tb(slide, M, H - Inches(0.62), W - 2 * M, Inches(0.3))
    _p(tf, text, 11.5, FAINT, first=True)


def new(kicker=None, title=None):
    s = prs.slides.add_slide(BLANK)
    top = header(s, title, kicker) if title else M
    return s, top


def divider(number, title, sub):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, ACCENT)
    tf = _tb(s, M, H / 2 - Inches(1.0), W - 2 * M, Inches(2.2))
    _p(tf, f"{number:02d}", 15, RGBColor(0x8F, 0xC5, 0xD4), bold=True, font=MONO, first=True)
    _p(tf, title, 40, WHITE, bold=True, space_before=4, space_after=10)
    _p(tf, sub, 17, RGBColor(0xC6, 0xDF, 0xE8))
    return s


# =========================================================================
# 1  title
# =========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, Inches(0.28), H, ACCENT)
tf = _tb(s, Inches(1.1), Inches(2.0), Inches(10.6), Inches(3.2))
_p(tf, "RENE / CUPDAQ", 15, ACCENT, bold=True, font=MONO, first=True)
_p(tf, "Text-mode Run Control", 46, INK, bold=True, space_before=6, space_after=2)
_p(tf, "rcterm · rcsupervisor · postrun", 46, INK, bold=True, space_after=16)
_p(tf, "A C++17 replacement for the rc.py GUI run control, with automated "
       "run rotation, fault recovery and post-run processing.", 18, MUTED)
tf2 = _tb(s, Inches(1.1), H - Inches(1.5), Inches(10.6), Inches(0.9))
_p(tf2, "Yeong-Gwang site   ·   Rocky Linux 9.8   ·   ROOT 6.28/04   ·   GCC 11.5.0",
   13, FAINT, font=MONO, first=True)
_p(tf2, "August 2026   ·   runs 4288-4290, two rotations verified", 13, FAINT, font=MONO)

# =========================================================================
# 2  what it is
# =========================================================================
s, top = new("overview", "Two programs, one job")
y = metric_row(s, [
    ("rcterm", "runs one run: boot → config → start → monitor → end", ACCENT),
    ("rcsupervisor", "rotates runs every N hours, diagnoses and recovers", ACCENT),
    ("postrun", "merges FADC+SADC and produces the analysis files", ACCENT2),
], top)
bullets(s, [
    ("The run control never touches the hardware. It speaks a 32-byte socket "
     "protocol to the CUPDAQ binaries — daq, merger, tcb — which do.", ""),
    ("No hardware driver was rewritten. What was rewritten is the orchestration: "
     "parsing the config, launching each node, sequencing the state machine, "
     "recording the run in the catalogue.", ""),
    ("rcsupervisor deliberately does not link the state machine. A bug in the run "
     "control must not take the watchdog down with it.", ""),
], y + Inches(0.45))

# =========================================================================
# 3  background
# =========================================================================
s, top = new("background", "Where this came from")
y = bullets(s, [
    ("rc.py — a PyQt5 GUI run control — required a desktop session to take data.",""),
    ("The August 14 incident made the cost concrete: the nouveau driver faulted, "
     "the screen froze, and the DAQ was running inside a terminal window of that "
     "very session.", ""),
    ("A text-mode run control detaches data taking from the desktop. It runs under "
     "tmux, survives the display dying, and can be driven over ssh.", ""),
], top)
codebox(s, [
    "22:37:02  nouveau GSP MMU fault → gnome-shell channel killed",
    "22:37:28  rcsupervisor recovered by itself, run 4276 started",
    "22:55:48  run ended cleanly by request — 18 min of data lost",
    "",
    "the DAQ survived the GPU fault. the session it ran in did not.",
], M, Inches(4.35), W - 2 * M, size=14)
footer(s, "Full account: docs kept in the repository under CLAUDE.md §11.2")

# =========================================================================
# 4  architecture
# =========================================================================
s, top = new("architecture", "How a run is assembled")
codebox(s, [
    "rcsupervisor            rotation timer, health checks, recovery",
    "     │  fork/exec (one run per cycle)",
    "     ▼",
    "  rcterm                config parse → state machine → run catalogue",
    "     │  executedaq.sh (one per node)      │  32-byte socket",
    "     ▼                                    ▼",
    "  daq -f   daq -s   tcb                 QUERY / CONFIG / START / END",
    "  (FADC)   (SADC)   (trigger)",
    "     │",
    "     ▼  raw files, one per subrun",
    "  FADC_%06d.root.%05d      SADC_%06d.root.%05d",
], M, top, W - 2 * M, size=14.5)
bullets(s, [
    ("Boot order is forced by node role — MERGER → ADC → TCB — not by name. "
     "Sorting by name puts AMOREADC last and the trigger board comes up after its "
     "clients, which breaks the connection.", ""),
], Inches(5.55), size=15)

# =========================================================================
# 5  protocol
# =========================================================================
s, top = new("architecture", "The protocol, established from the source")
y = bullets(s, [
    ("Every message is 32 bytes: four little-endian 8-byte unsigned values.", ""),
    ("State is a bitmask, not an integer. status & (1 << state).", ""),
    ("This matters: a leftover run reports status 0x8 — Running, with no Booted "
     "bit — and a naive integer comparison reads that as a valid state.", ""),
], top)
codebox(s, [
    "CONFIGRUN 1     STARTRUN 2      ENDRUN 3        EXIT 4",
    "QUERYDAQSTATUS 10   QUERYRUNINFO 12   QUERYTRGINFO 14   QUERYMONITOR 21",
    "",
    "Down 0   Booted 1   Configured 2   Running 3",
    "RunEnded 4   ProcEnded 5   Warning 6   Error 7",
], M, y + Inches(0.2), W - 2 * M, size=14.5)

# =========================================================================
# 6  divider - verification
# =========================================================================
divider(1, "Proving it works without a detector",
        "Hardware time is scarce. Almost everything was verified before touching it.")

# =========================================================================
# 7  fake TCB
# =========================================================================
s, top = new("verification", "A fake detector that speaks the real protocol")
y = bullets(s, [
    ("A 120-line stand-in answers the 32-byte protocol on the TCB and ADC ports, "
     "and a fake executedaq.sh stands in for the launcher. No hardware is touched.", ""),
    ("It can be told to misbehave — to delay the RunEnded transition, or to stop "
     "answering — which is how the failure paths get exercised at all.", ""),
], top)
table(s, [
    ["Test", "Result"],
    ["Rotation: run ends, next run starts", "+clean exit, no zombies"],
    ["Stale heartbeat detected and recovered", "+detected at 12 s, restarted"],
    ["Graceful stop finalises the run", "+exit 0, catalogue complete"],
    ["Same case on the pre-fix binary", "-reproduced the field symptom"],
    ["Boot failure is labelled in the catalogue", "+onlbit=0 with a reason"],
    ["Second signal escapes a stuck DAQ", "+exited within 1 s"],
], y + Inches(0.25), [Inches(7.4), Inches(4.4)])

# =========================================================================
# 8  divider - defects
# =========================================================================
divider(2, "Three defects found in operation",
        "All three were invisible until the system ran unattended against real data.")

# =========================================================================
# 9  defect 1
# =========================================================================
s, top = new("defect 1 · the serious one", "Run finalisation lost on every rotation")
y = bullets(s, [
    ("WaitState returned false the moment a stop was requested. On SIGTERM — which "
     "is how the supervisor ends every run — rcterm sent ENDRUN and then abandoned "
     "the wait, so FinalizeRunInDB was never reached.", ""),
    ("The data files were complete. The catalogue row was not: no start time, no end "
     "time, no event counts. The heartbeat said failed and the exit code said failure, "
     "for a run that had ended perfectly.", ""),
], top)
codebox(s, [
    "healthy run   ENDRUN → ENDED → EXIT",
    "runs 4276,    ENDRUN →          EXIT          ← no ENDED. row left empty",
    "4284, 4287",
], M, y + Inches(0.15), W - 2 * M, size=14.5)
bullets(s, [
    ("The end-of-run confirmation now ignores the stop request; boot and configure "
     "still abort at once. A second signal means leave now, so a stuck DAQ can never "
     "trap the operator.", ""),
], Inches(5.45), size=15)

# =========================================================================
# 10  defect 1 evidence
# =========================================================================
s, top = new("defect 1 · evidence", "Proven by rebuilding the previous version")
bullets(s, [
    ("A binary built from the commit before the fix was run against the same scenario "
     "as the fixed one. It reproduced the field symptom exactly — which is what "
     "makes this a diagnosis rather than a guess.", ""),
], top, size=16)
y = table(s, [
    ["Same scenario, same fake detector", "Before", "After"],
    ["Heartbeat phase", "-failed", "+ended"],
    ["Exit code", "-2  (failure)", "+0  (success)"],
    ["Catalogue start / end time", "-empty", "+recorded"],
    ["Per-ADC event counts", "-empty", "+15811 / 15811"],
], top + Inches(1.05), [Inches(6.0), Inches(2.9), Inches(2.9)])
codebox(s, [
    "field evidence   35 catalogue rows left unusable across runs 4247-4284",
    "                 every supervisor rotation hit this path",
], M, y + Inches(0.35), W - 2 * M, size=14)

# =========================================================================
# 11  defect 2
# =========================================================================
s, top = new("defect 2", "Attaching to the previous run's trigger board")
y = bullets(s, [
    ("If a DAQ from an earlier run was still alive, the newly launched tcb could not "
     "bind the port and died. rcterm then connected to that port anyway and drove the "
     "old run's trigger board.", ""),
    ("It reported status 0x8 — Running, never Booted — so rcterm waited out the "
     "whole boot timeout and failed, burning a run number each time.", ""),
], top)
codebox(s, [
    "TCB_004269.log",
    "  22:23:55  new client connection      ← this is run 4270's rcterm,",
    "                                          talking to run 4269's trigger board",
], M, y + Inches(0.15), W - 2 * M, size=14.5)
bullets(s, [
    ("The port is now checked before the run number is issued, so a stale DAQ costs "
     "no number, and the message names the tool that clears it.", ""),
], Inches(5.25), size=15)

# =========================================================================
# 12  defect 3
# =========================================================================
s, top = new("defect 3", "Failed runs left no trace of why")
y = bullets(s, [
    ("A run number is issued before booting. If the boot fails, the row stays behind "
     "with empty times and nothing to say what happened.", ""),
    ("After a spell of repeated failures the catalogue holds a block of identical "
     "empty rows, and a run whose record is merely missing looks exactly like one "
     "that never started.", ""),
], top)
codebox(s, [
    "onlbit = 0   runlog = 'boot failed; run never started'",
    "onlbit = 0   runlog = 'aborted; run started but was not finalized'",
], M, y + Inches(0.2), W - 2 * M, size=15)
bullets(s, [
    ("The two cases are now distinguished by whether the run ever reached STARTRUN.", ""),
], Inches(4.85), size=15)

# =========================================================================
# 13  divider - pipeline
# =========================================================================
divider(3, "Post-run processing, attached to the DAQ",
        "Merging and production now follow acquisition instead of waiting for it.")

# =========================================================================
# 14  pipeline flow
# =========================================================================
s, top = new("pipeline", "From raw subruns to analysis files")
codebox(s, [
    "FADC_%06d.root.%05d   (target, ~70 MB)   ┐",
    "                                          ├─ merge_FADC_SADC_v3_5v.cc",
    "SADC_%06d.root.%05d   (VETO,   ~8 MB)    ┘        │   match by trigger number",
    "                                                   ▼",
    "                        Merged/MERGED_%06d.root.%05d   (~80 MB)  + QC canvas",
    "                                                   │",
    "                                production_from_merged_v3_5v.cc",
    "                                                   ▼",
    "                        PRD/PRD_%06d.%05d.root          (~77 MB)",
], M, top, W - 2 * M, size=14)
bullets(s, [
    ("The physics code is called where it lives and is not copied into this project. "
     "A copy would fork, and the two would drift apart.", ""),
    ("Completion is judged the same way as before, and the log file names are "
     "unchanged, so the existing audit_run.sh still works on this output.", ""),
], Inches(4.85), size=15.5)

# =========================================================================
# 15  why only production parallel
# =========================================================================
s, top = new("pipeline", "Only one of the two stages can be parallelised")
y = table(s, [
    ["Stage", "Time per subrun", "Parallel?"],
    ["merge", "28 s", "-no"],
    ["production", "15 s", "+yes"],
    ["total, serial", "43 s", ""],
], top, [Inches(5.2), Inches(3.5), Inches(3.1)])
bullets(s, [
    ("Measured over 33,357 subrun logs from runs 4238, 4239 and 4240.", "sub"),
    ("A subrun holds 60 s of data, so the serial pipeline already keeps up — but "
     "with only 28% to spare, and contention has been observed to double it.", ""),
    ("merge cannot be split: the macro prints the SADC position at the end of each "
     "subrun and the next subrun has to start from it, or events are lost at the "
     "boundary. production is independent per subrun, so it is run in a pool.", ""),
], y + Inches(0.3), size=16)

# =========================================================================
# 16  trailing
# =========================================================================
s, top = new("pipeline", "Trailing the acquisition by three minutes")
y = bullets(s, [
    ("A 24 h run holds about 1440 subruns. Processing only after the run ends would "
     "leave a 17-hour tail every day, so processing follows acquisition instead.", ""),
    ("The safe edge comes from rcterm's own heartbeat, whose subrun field is the file "
     "currently being written.", ""),
], top, size=16)
codebox(s, [
    "subrun 845   being written   ← never touched",
    "subrun 844   complete        ┐",
    "subrun 843   complete        ┘ margin: flush, and clock skew",
    "subrun 842   ────→  processing stops here   (heartbeat subrun - 3)",
], M, y + Inches(0.2), W - 2 * M, size=14.5)
bullets(s, [
    ("The margin is not decoration. The NFS server clock runs about 28 s ahead of the "
     "local one, so file timestamps alone cannot establish that a file is finished; "
     "opening one too early makes the merge declare it corrupt and retry.", ""),
], Inches(5.3), size=15)

# =========================================================================
# 17  bottleneck
# =========================================================================
s, top = new("pipeline · measurement", "The bottleneck is storage, not CPU")
y = metric_row(s, [
    ("26-34%", "iowait while processing", BAD),
    ("61-88%", "CPU idle at the same time", MUTED),
    ("0.14", "cores used by the DAQ itself, of 12", MUTED),
], top)
bullets(s, [
    ("Adding parallel jobs therefore buys almost nothing — it only adds contention "
     "on the same filesystem the DAQ is writing to.", ""),
], y + Inches(0.3), size=16)
y2 = table(s, [
    ["Where the merged and produced files are written", "Per subrun", "Margin vs 60 s"],
    ["/scratch  (NFS)", "41.0 s", "1.5x"],
    ["/Data_ssd  (local NVMe)", "+27.7 s", "+2.1x"],
], y + Inches(0.95), [Inches(6.6), Inches(2.6), Inches(2.6)])
footer(s, "Raw data stays on NFS; only the derived files move. The macros are unchanged — the directories are symlinked.")

# =========================================================================
# 18  divider - operations
# =========================================================================
divider(4, "Running it", "What an operator actually sees and types.")

# =========================================================================
# 19  screen
# =========================================================================
s, top = new("operations", "One screen, four panes")
codebox(s, [
    "+--------------------+-------------------------+",
    "| monitor         6  |                         |   left   DAQ status",
    "+--------------------+                         |   right  your shell",
    "| supervisor      2  |    work space           |",
    "+--------------------+                         |   survives ssh drops,",
    "| postrun         2  |                         |   display crashes,",
    "+--------------------+-------------------------+   terminal restarts",
    "        4.5                     5.5",
], M, top, W - 2 * M, size=14)
bullets(s, [
    ("The supervisor forces rcterm into one-line mode, because the full screen clears "
     "itself and would erase the supervisor's own messages. The monitor pane rebuilds "
     "that screen from the heartbeat file instead — read-only, so it can be opened "
     "as many times as wanted, and it shows the heartbeat age, which makes a stalled "
     "run visible rather than merely frozen.", ""),
], Inches(4.55), size=15.5)

# =========================================================================
# 20  commands
# =========================================================================
s, top = new("operations", "The commands that matter")
codebox(s, [
    "# bring up the whole screen",
    "scripts/daq-tmux.sh                 # layout only, touches no hardware",
    "scripts/daq-tmux.sh --start         # also starts the DAQ if it is down",
    "",
    "# re-attach later, from anywhere",
    "tmux attach -t daq                  # Ctrl-B then D to detach",
    "Ctrl-B  =                           # restore the pane proportions",
    "",
    "# stop the DAQ safely  —  the supervisor PID, and only that one",
    "kill -TERM $(pgrep -x rcsupervisor)",
    "",
    "# post-run processing",
    "scripts/postrun.sh --once --dry-run           # what would it do",
    "scripts/postrun.sh --follow --outroot /Data_ssd/RAW",
], M, top, W - 2 * M, size=14)
bullets(s, [
    ("Never signal the process group. daq, tcb and merger are not children of rcterm; "
     "killing the group cuts them off mid-write and corrupts the last file.", ""),
], Inches(5.75), size=15)

# =========================================================================
# 21  storage
# =========================================================================
s, top = new("operations", "What a day of running costs")
y = table(s, [
    ["Per subrun (60 s of data)", "Size"],
    ["raw  FADC + SADC", "78 MB"],
    ["merged", "80 MB"],
    ["produced", "77 MB"],
    ["total", "235 MB"],
], top, [Inches(6.0), Inches(3.0)])
metric_row(s, [
    ("334 GB", "per 24-hour run", ACCENT2),
    ("48 days", "headroom on /scratch", MUTED),
    ("9 days", "headroom on /Data_ssd", BAD),
], y + Inches(0.4))
bullets(s, [
    ("The local disk is the faster place to write and the smaller place to keep. "
     "A retention policy is needed before it fills, not after.", ""),
], y + Inches(2.2), size=16)

# =========================================================================
# 22  status
# =========================================================================
s, top = new("status", "Where things stand")
y = bullets(s, [
    ("Two 24-hour rotations have now passed under the fix. Both logged ENDED, "
     "exited 0, and wrote the catalogue row in full. One could be luck; two is a "
     "reproduction.", "good"),
    ("A sweep over the day and a half found nothing to fix: no recovery or FATAL "
     "entries against 286 passing health checks, no zombie subruns, no orphan rows.", "good"),
    ("Both runs came out even at every stage — 1440 FADC, 1440 SADC, 1440 merged, "
     "1440 produced — so nothing was dropped on either arm.", "good"),
    ("Post-run processing keeps pace with acquisition, three subruns behind.", "good"),
], top, size=16)
bullets(s, [
    ("The open item is storage, not correctness", "head"),
    ("Each run leaves about 217 GB of derived files on the local disk and 1.5 TB "
     "is free — roughly a week. Finished runs have to move back to NFS, and the "
     "symlinks make that fiddly enough by hand that it belongs in the driver.", ""),
], y + Inches(0.2), size=16)

# =========================================================================
# 23  close
# =========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, PANEL)
rect(s, 0, 0, Inches(0.28), H, ACCENT)
tf = _tb(s, Inches(1.2), Inches(1.5), Inches(10.8), Inches(4.5))
_p(tf, "WHAT CHANGED, IN ONE LINE", 13, ACCENT, bold=True, font=MONO, first=True)
_p(tf, "Data taking no longer depends on a desktop session,",
   30, INK, bold=True, space_before=14, space_after=2)
_p(tf, "a rotation no longer loses the run record,", 30, INK, bold=True, space_after=2)
_p(tf, "and analysis files are ready when the run ends.", 30, INK, bold=True, space_after=22)
_p(tf, "github.com/Sang-Yong/RENE-daq-rcterm", 16, MUTED, font=MONO)
_p(tf, "docs/MANUAL.md   ·   docs/POSTRUN.md   ·   config/dotfiles/README.md",
   16, MUTED, font=MONO, space_before=4)

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "deck.pptx"
prs.save(out)
print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")
