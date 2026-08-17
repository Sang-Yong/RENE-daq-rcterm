#!/usr/bin/env python3
"""RENE DAQ run control - presentation generator."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette -------------------------------------------------------------
INK      = RGBColor(0x16, 0x1B, 0x22)
MUTED    = RGBColor(0x5A, 0x66, 0x72)
FAINT    = RGBColor(0x8C, 0x98, 0xA4)
ACCENT   = RGBColor(0x0E, 0x5C, 0x73)   # deep teal
ACCENT2  = RGBColor(0xC2, 0x5E, 0x00)   # burnt orange
GOOD     = RGBColor(0x1B, 0x6B, 0x45)
BAD      = RGBColor(0xA3, 0x28, 0x20)
RULE     = RGBColor(0xD8, 0xDF, 0xE4)
PANEL    = RGBColor(0xF2, 0xF5, 0xF7)
CODEBG   = RGBColor(0x16, 0x1B, 0x22)
CODEFG   = RGBColor(0xD6, 0xDE, 0xE4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SANS = "맑은 고딕"
MONO = "D2Coding"          # 없으면 PowerPoint 가 대체한다
SANS_EA = "맑은 고딕"
MONO_EA = "맑은 고딕"


def _set_font(run, name):
    """latin 뿐 아니라 ea(한중일) / cs 타이프페이스까지 지정한다.
    이걸 안 하면 한글이 PowerPoint 기본 글꼴로 제멋대로 바뀐다."""
    from pptx.oxml.ns import qn
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = name if name != MONO else MONO_EA
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", ea)

W, H = Inches(13.333), Inches(7.5)
M    = Inches(0.75)                      # margin

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _p(tf, text, size, color, bold=False, font=SANS, space_before=0,
       space_after=4, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    r.font.size, r.font.bold = Pt(size), bold
    _set_font(r, font)
    r.font.color.rgb = color
    return p


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def header(slide, title, kicker=None):
    y = M
    if kicker:
        tf = _tb(slide, M, y - Inches(0.06), W - 2 * M, Inches(0.3))
        _p(tf, kicker.upper(), 12, ACCENT, bold=True, font=MONO, first=True)
        y += Inches(0.34)
    tf = _tb(slide, M, y, W - 2 * M, Inches(0.8))
    _p(tf, title, 30, INK, bold=True, first=True)
    rect(slide, M, y + Inches(0.72), Inches(1.1), Pt(3), ACCENT)
    return y + Inches(1.05)


def bullets(slide, items, top, size=17, width=None, left=None, gap=9):
    """items: list of (text, kind) where kind in {'', 'sub', 'good', 'bad', 'code'}"""
    left = left or M
    width = width or (W - 2 * M)
    tf = _tb(slide, left, top, width, H - top - M)
    # 대략적인 줄 수로 아래쪽 y 를 되돌려 준다. 다음 요소를 이어 붙이기 위한 것.
    chars_per_line = max(20, int(width / Inches(1.0) * 11.6 * (17.0 / size)))
    used = Emu(0)
    first = True

    def _w(t):
        # 한글은 라틴 문자보다 넓다. 줄 수 추정에서 1.8 배로 센다.
        return sum(1.8 if '\uac00' <= c <= '\ud7a3' or '\u3130' <= c <= '\u318f'
                   else 1.0 for c in t)

    for text, kind in items:
        nlines = max(1, -(-int(_w(text)) // chars_per_line))
        fs = size - 3 if kind in ("sub", "code") else size
        used += Emu(int(Pt(fs * 1.28) * nlines + Pt(gap)))
    for text, kind in items:
        if kind == "sub":
            _p(tf, "     " + text, size - 3, MUTED, space_after=gap - 2, first=first)
        elif kind == "good":
            _p(tf, "✓  " + text, size, GOOD, bold=True, space_after=gap, first=first)
        elif kind == "bad":
            _p(tf, "✗  " + text, size, BAD, bold=True, space_after=gap, first=first)
        elif kind == "code":
            _p(tf, text, size - 3, ACCENT, font=MONO, space_after=gap, first=first)
        elif kind == "head":
            _p(tf, text, size, ACCENT, bold=True, space_before=8, space_after=gap, first=first)
        else:
            _p(tf, "•  " + text, size, INK, space_after=gap, first=first)
        first = False
    return top + used


def codebox(slide, lines, x, y, w, h=None, size=13, title=None):
    lh = Pt(size * 1.45)
    h = h or Emu(int(lh * (len(lines) + 1.4)))
    rect(slide, x, y, w, h, CODEBG)
    tf = _tb(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), h - Inches(0.3))
    first = True
    if title:
        _p(tf, title, size - 1, FAINT, font=MONO, space_after=5, first=True)
        first = False
    for ln in lines:
        col = CODEFG
        if ln.startswith("#") or ln.startswith("//"):
            col = FAINT
        _p(tf, ln, size, col, font=MONO, space_after=1, first=first)
        first = False
    return y + h


def metric_row(slide, items, y, h=Inches(1.5)):
    """items: list of (big, label, color)"""
    n = len(items)
    gap = Inches(0.25)
    w = (W - 2 * M - gap * (n - 1)) / n
    for i, (big, label, col) in enumerate(items):
        x = M + i * (w + gap)
        rect(slide, x, y, w, h, PANEL)
        rect(slide, x, y, Pt(4), h, col)
        tf = _tb(slide, x + Inches(0.28), y + Inches(0.22), w - Inches(0.5), h - Inches(0.4))
        _p(tf, big, 34, col, bold=True, space_after=2, first=True)
        _p(tf, label, 12.5, MUTED)
    return y + h


def table(slide, rows, top, widths, size=14, head=True):
    x = M
    y = top
    rowh = Inches(0.42)
    for ri, row in enumerate(rows):
        cx = x
        if ri == 0 and head:
            rect(slide, M, y, sum(widths), rowh, PANEL)
        for ci, cell in enumerate(row):
            tf = _tb(slide, cx + Inches(0.14), y + Inches(0.08), widths[ci] - Inches(0.2), rowh)
            col = MUTED if (ri == 0 and head) else INK
            bold = (ri == 0 and head)
            fnt = SANS
            if cell.startswith("`") and cell.endswith("`"):
                cell = cell[1:-1]; fnt = MONO
            if cell.startswith("+"):
                col = GOOD; cell = cell[1:]
            elif cell.startswith("-"):
                col = BAD; cell = cell[1:]
            _p(tf, cell, size, col, bold=bold, font=fnt, first=True)
            cx += widths[ci]
        y += rowh
        if ri == 0 and head:
            rect(slide, M, y, sum(widths), Pt(1), RULE)
        elif ri < len(rows) - 1:
            rect(slide, M, y, sum(widths), Pt(0.75), RULE)
    return y


def footer(slide, text):
    tf = _tb(slide, M, H - Inches(0.62), W - 2 * M, Inches(0.3))
    _p(tf, text, 11.5, FAINT, first=True)


def new(kicker=None, title=None):
    s = prs.slides.add_slide(BLANK)
    top = header(s, title, kicker) if title else M
    return s, top


def divider(number, title, sub):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, ACCENT)
    tf = _tb(s, M, H / 2 - Inches(1.0), W - 2 * M, Inches(2.2))
    _p(tf, f"{number:02d}", 15, RGBColor(0x8F, 0xC5, 0xD4), bold=True, font=MONO, first=True)
    _p(tf, title, 40, WHITE, bold=True, space_before=4, space_after=10)
    _p(tf, sub, 17, RGBColor(0xC6, 0xDF, 0xE8))
    return s


# =========================================================================
# 1  title
# =========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, Inches(0.28), H, ACCENT)
tf = _tb(s, Inches(1.1), Inches(2.0), Inches(10.6), Inches(3.2))
_p(tf, "RENE / CUPDAQ", 15, ACCENT, bold=True, font=MONO, first=True)
_p(tf, "텍스트 런컨트롤", 46, INK, bold=True, space_before=6, space_after=2)
_p(tf, "rcterm · rcsupervisor · postrun", 46, INK, bold=True, space_after=16)
_p(tf, "rc.py GUI 런컨트롤을 C++17 로 다시 쓴 것. 런 자동 로테이션, "
       "장애 자동 복구, 수집 뒤 후처리까지 포함한다.", 18, MUTED)
tf2 = _tb(s, Inches(1.1), H - Inches(1.5), Inches(10.6), Inches(0.9))
_p(tf2, "Yeong-Gwang site   ·   Rocky Linux 9.8   ·   ROOT 6.28/04   ·   GCC 11.5.0",
   13, FAINT, font=MONO, first=True)
_p(tf2, "2026년 8월   ·   run 4288~4290, 로테이션 2회 검증", 13, FAINT, font=MONO)

# =========================================================================
# 2  what it is
# =========================================================================
s, top = new("개요", "프로그램 셋, 하나의 일")
y = metric_row(s, [
    ("rcterm", "런 1회 : boot → config → start → 감시 → end", ACCENT),
    ("rcsupervisor", "N시간마다 런 교체, 이상 진단과 자동 복구", ACCENT),
    ("postrun", "FADC+SADC 병합과 분석용 PRD 생성", ACCENT2),
], top)
bullets(s, [
    ("런컨트롤은 하드웨어를 직접 건드리지 않는다. 32바이트 소켓 프로토콜로 "
     "CUPDAQ 바이너리(daq, merger, tcb)에 지시할 뿐이다.", ""),
    ("하드웨어 드라이버는 한 줄도 다시 쓰지 않았다. 다시 쓴 것은 조율이다 — "
     "config 파싱, 노드별 기동, 상태머신 순서 제어, 런 카탈로그 기록.", ""),
    ("rcsupervisor 는 상태머신을 일부러 링크하지 않는다. 런컨트롤의 버그가 "
     "감시자까지 끌고 내려가면 자동 복구라는 목적 자체가 무너진다.", ""),
], y + Inches(0.45))

# =========================================================================
# 3  background
# =========================================================================
s, top = new("배경", "왜 만들었나")
y = bullets(s, [
    ("rc.py 는 PyQt5 GUI 런컨트롤이라 데이터를 받으려면 데스크톱 세션이 필요했다.",""),
    ("8월 14일 사고로 그 대가가 분명해졌다. nouveau 드라이버가 fault 를 내 화면이 "
     "멎었는데, DAQ 가 하필 그 세션의 터미널 창 안에서 돌고 있었다.", ""),
    ("텍스트 런컨트롤은 수집을 데스크톱에서 떼어낸다. tmux 안에서 돌고, 화면이 "
     "죽어도 살아남고, ssh 로 조작할 수 있다.", ""),
], top)
codebox(s, [
    "22:37:02  nouveau GSP MMU fault → gnome-shell 채널 kill",
    "22:37:28  rcsupervisor 가 스스로 복구, run 4276 시작",
    "22:55:48  요청대로 정상 종료 — 18분 손실",
    "",
    "DAQ 는 GPU fault 를 견뎠다. 그것이 돌던 세션이 못 견뎠다.",
], M, Inches(4.35), W - 2 * M, size=14)
footer(s, "전체 기록은 저장소의 CLAUDE.md §11.2 에 있다")

# =========================================================================
# 4  architecture
# =========================================================================
s, top = new("구조", "런 하나가 조립되는 과정")
codebox(s, [
    "rcsupervisor            로테이션 타이머, 상태 점검, 복구",
    "     │  fork/exec (사이클당 런 1개)",
    "     ▼",
    "  rcterm                config 파싱 → 상태머신 → 런 카탈로그",
    "     │  executedaq.sh (노드마다 하나)      │  32바이트 소켓",
    "     ▼                                    ▼",
    "  daq -f   daq -s   tcb                 QUERY / CONFIG / START / END",
    "  (FADC)   (SADC)   (트리거)",
    "     │",
    "     ▼  서브런마다 원시 파일",
    "  FADC_%06d.root.%05d      SADC_%06d.root.%05d",
], M, top, W - 2 * M, size=14.5)
bullets(s, [
    ("부팅 순서는 이름이 아니라 노드 역할로 강제한다 — MERGER → ADC → TCB. "
     "이름으로 정렬하면 AMOREADC 가 마지막이 되어 트리거 보드가 클라이언트보다 "
     "늦게 떠 접속이 깨진다.", ""),
], Inches(5.55), size=15)

# =========================================================================
# 5  protocol
# =========================================================================
s, top = new("구조", "프로토콜 — 소스 실측으로 확정")
y = bullets(s, [
    ("모든 메시지는 32바이트 — 리틀엔디언 8바이트 부호없는 정수 4개.", ""),
    ("상태는 정수가 아니라 비트마스크다. status & (1 << state).", ""),
    ("이게 중요하다. 남아있던 런은 status 0x8 을 답한다 — Running 인데 Booted "
     "비트가 없다. 정수로 비교하면 정상 상태로 오해한다.", ""),
], top)
codebox(s, [
    "CONFIGRUN 1     STARTRUN 2      ENDRUN 3        EXIT 4",
    "QUERYDAQSTATUS 10   QUERYRUNINFO 12   QUERYTRGINFO 14   QUERYMONITOR 21",
    "",
    "Down 0   Booted 1   Configured 2   Running 3",
    "RunEnded 4   ProcEnded 5   Warning 6   Error 7",
], M, y + Inches(0.2), W - 2 * M, size=14.5)

# =========================================================================
# 6  divider - verification
# =========================================================================
divider(1, "검출기 없이 증명하기",
        "하드웨어 시간은 귀하다. 거의 전부를 건드리기 전에 검증했다.")

# =========================================================================
# 7  fake TCB
# =========================================================================
s, top = new("검증", "진짜 프로토콜을 말하는 가짜 검출기")
y = bullets(s, [
    ("120줄짜리 대역이 TCB·ADC 포트에서 32바이트 프로토콜에 답하고, 가짜 "
     "executedaq.sh 가 기동 스크립트를 대신한다. 하드웨어는 전혀 건드리지 않는다.", ""),
    ("일부러 오작동시킬 수도 있다 — RunEnded 전이를 늦추거나, 아예 응답을 끊거나. "
     "실패 경로는 그렇게 해야만 시험된다.", ""),
], top)
table(s, [
    ["시험", "결과"],
    ["로테이션 : 런 종료 후 다음 런 시작", "+정상 종료, 좀비 없음"],
    ["heartbeat 정지 감지와 복구", "+12초에 감지, 재시작"],
    ["정상 정지 시 런 마감", "+exit 0, 카탈로그 완전"],
    ["수정 전 바이너리로 같은 시나리오", "-현장 증상 그대로 재현"],
    ["부팅 실패가 카탈로그에 표기", "+onlbit=0 과 사유"],
    ["멈춘 DAQ 에서 2차 신호로 탈출", "+1초 내 종료"],
], y + Inches(0.25), [Inches(7.4), Inches(4.4)])

# =========================================================================
# 8  divider - defects
# =========================================================================
divider(2, "실운용에서 찾은 결함 3종",
        "셋 다 실제 데이터로 무인 운용하기 전에는 드러나지 않았다.")

# =========================================================================
# 9  defect 1
# =========================================================================
s, top = new("결함 1 · 가장 심각", "로테이션마다 런 기록이 유실됐다")
y = bullets(s, [
    ("정지 요청이 오는 순간 WaitState 가 false 를 반환했다. 감시자는 모든 런을 "
     "SIGTERM 으로 끝내는데, rcterm 은 ENDRUN 만 보내고 대기를 포기해 "
     "FinalizeRunInDB 에 도달하지 못했다.", ""),
    ("데이터 파일은 멀쩡했다. 카탈로그 행이 비었다 — 시작·종료 시각도, 이벤트 "
     "수도 없다. 완벽하게 끝난 런인데 heartbeat 는 failed, 종료코드는 실패였다.", ""),
], top)
codebox(s, [
    "정상 런        ENDRUN → ENDED → EXIT",
    "run 4276,     ENDRUN →          EXIT      ← ENDED 없음. 행이 빈 채로 남음",
    "4284, 4287",
], M, y + Inches(0.15), W - 2 * M, size=14.5)
bullets(s, [
    ("이제 런 종료 확인만 정지 요청을 무시한다. 부팅과 설정 단계는 그대로 즉시 "
     "중단한다. 두 번째 신호는 즉시 나가라는 뜻이라, 멈춘 DAQ 가 운영자를 "
     "가둘 일은 없다.", ""),
], Inches(5.45), size=15)

# =========================================================================
# 10  defect 1 evidence
# =========================================================================
s, top = new("결함 1 · 증거", "직전 버전을 다시 빌드해 증명했다")
bullets(s, [
    ("수정 직전 커밋으로 빌드한 바이너리를 수정본과 같은 시나리오에 돌렸다. "
     "현장 증상을 그대로 재현했다. 이것이 추측이 아니라 진단인 이유다.", ""),
], top, size=16)
y = table(s, [
    ["같은 시나리오, 같은 가짜 검출기", "수정 전", "수정 후"],
    ["heartbeat phase", "-failed", "+ended"],
    ["종료코드", "-2  (실패)", "+0  (성공)"],
    ["카탈로그 시작/종료 시각", "-빈 값", "+기록됨"],
    ["ADC 별 이벤트 수", "-빈 값", "+15811 / 15811"],
], top + Inches(1.05), [Inches(6.0), Inches(2.9), Inches(2.9)])
codebox(s, [
    "현장 증거   run 4247~4284 구간에서 쓸 수 없게 된 카탈로그 행 35개",
    "            감시자 로테이션은 매번 이 경로를 지났다",
], M, y + Inches(0.35), W - 2 * M, size=14)

# =========================================================================
# 11  defect 2
# =========================================================================
s, top = new("결함 2", "이전 런의 트리거 보드에 접속했다")
y = bullets(s, [
    ("이전 런의 DAQ 가 살아있으면 새로 띄운 tcb 가 포트를 잡지 못하고 죽는다. "
     "그런데 rcterm 은 그 포트에 그냥 접속해 옛 런의 트리거 보드를 조작했다.", ""),
    ("그 보드는 status 0x8 을 답한다 — Running 인데 Booted 가 없다. rcterm 은 "
     "부팅 타임아웃을 다 기다렸다 실패하고, 그때마다 run 번호를 하나씩 버렸다.", ""),
], top)
codebox(s, [
    "TCB_004269.log",
    "  22:23:55  새 클라이언트 접속      ← run 4270 의 rcterm 이",
    "                                       run 4269 의 트리거 보드와 통신 중",
], M, y + Inches(0.15), W - 2 * M, size=14.5)
bullets(s, [
    ("이제 run 번호를 발급하기 전에 포트를 확인한다. 남은 DAQ 때문에 번호를 "
     "버리지 않고, 메시지가 정리 도구 이름까지 알려준다.", ""),
], Inches(5.25), size=15)

# =========================================================================
# 12  defect 3
# =========================================================================
s, top = new("결함 3", "실패한 런이 이유를 남기지 않았다")
y = bullets(s, [
    ("run 번호는 부팅 전에 발급된다. 부팅이 실패하면 시각이 빈 행만 남고 "
     "무슨 일이 있었는지 아무 것도 남지 않는다.", ""),
    ("반복 실패하면 카탈로그에 똑같이 빈 행이 뭉텅이로 쌓인다. 기록만 누락된 "
     "런과 아예 시작도 못 한 런이 구분되지 않는다.", ""),
], top)
codebox(s, [
    "onlbit = 0   runlog = 'boot failed; run never started'",
    "onlbit = 0   runlog = 'aborted; run started but was not finalized'",
], M, y + Inches(0.2), W - 2 * M, size=15)
bullets(s, [
    ("이제 STARTRUN 까지 갔는지로 두 경우를 구분한다.", ""),
], Inches(4.85), size=15)

# =========================================================================
# 13  divider - pipeline
# =========================================================================
divider(3, "DAQ 뒤에 붙인 후처리",
        "병합과 production 이 수집이 끝나기를 기다리지 않고 뒤따라간다.")

# =========================================================================
# 14  pipeline flow
# =========================================================================
s, top = new("파이프라인", "원시 서브런에서 분석 파일까지")
codebox(s, [
    "FADC_%06d.root.%05d   (타겟, ~70 MB)    ┐",
    "                                          ├─ merge_FADC_SADC_v3_5v.cc",
    "SADC_%06d.root.%05d   (VETO,  ~8 MB)    ┘        │   트리거 번호로 정합",
    "                                                   ▼",
    "                        Merged/MERGED_%06d.root.%05d   (~80 MB)  + QC 그림",
    "                                                   │",
    "                                production_from_merged_v3_5v.cc",
    "                                                   ▼",
    "                        PRD/PRD_%06d.%05d.root          (~77 MB)",
], M, top, W - 2 * M, size=14)
bullets(s, [
    ("물리 코드는 있는 자리에서 호출할 뿐 이 프로젝트로 복제하지 않는다. "
     "복제하면 두 벌이 갈라진다.", ""),
    ("완료 판정 방식과 로그 파일명을 원본 그대로 유지했다. 그래서 기존 "
     "audit_run.sh 가 이 결과에도 그대로 동작한다.", ""),
], Inches(4.85), size=15.5)

# =========================================================================
# 15  why only production parallel
# =========================================================================
s, top = new("파이프라인", "두 단계 중 하나만 병렬화할 수 있다")
y = table(s, [
    ["단계", "서브런당 소요", "병렬화"],
    ["merge", "28 초", "-불가"],
    ["production", "15 초", "+가능"],
    ["직렬 합계", "43 초", ""],
], top, [Inches(5.2), Inches(3.5), Inches(3.1)])
bullets(s, [
    ("run 4238 / 4239 / 4240 의 서브런 로그 33,357개 실측.", "sub"),
    ("서브런 1개는 60초 분량이라 직렬로도 따라간다 — 다만 여유가 28% 뿐이고, "
     "경합이 생기면 두 배까지 늘어지는 것을 실측했다.", ""),
    ("merge 는 나눌 수 없다. 매크로가 서브런 끝에 찍는 SADC 위치를 다음 서브런이 "
     "받아야 하고, 안 받으면 경계에서 이벤트를 잃는다. production 은 서브런마다 "
     "독립이라 풀로 돌린다.", ""),
], y + Inches(0.3), size=16)

# =========================================================================
# 16  trailing
# =========================================================================
s, top = new("파이프라인", "수집보다 3분 뒤에서 따라간다")
y = bullets(s, [
    ("24시간 런은 서브런이 약 1440개다. 런이 끝난 뒤에 몰아서 처리하면 매일 "
     "17시간짜리 꼬리가 남는다. 그래서 수집을 뒤따라가며 처리한다.", ""),
    ("안전선은 rcterm 자신의 heartbeat 에서 얻는다. 그 subrun 값이 바로 지금 "
     "기록 중인 파일 번호다.", ""),
], top, size=16)
codebox(s, [
    "subrun 845   기록 중       ← 절대 건드리지 않는다",
    "subrun 844   완료          ┐",
    "subrun 843   완료          ┘ 여유 : flush 와 시계 오차 흡수",
    "subrun 842   ────→  처리 상한         (heartbeat subrun - 3)",
], M, y + Inches(0.2), W - 2 * M, size=14.5)
bullets(s, [
    ("이 여유는 장식이 아니다. NFS 서버 시계가 로컬보다 약 28초 앞서 있어 "
     "파일 시각만으로는 완료를 판정할 수 없다. 덜 쓰인 파일을 열면 merge 가 "
     "손상으로 판정해 재시도에 들어간다.", ""),
], Inches(5.3), size=15)

# =========================================================================
# 17  bottleneck
# =========================================================================
s, top = new("파이프라인 · 실측", "병목은 CPU 가 아니라 저장장치다")
y = metric_row(s, [
    ("26~34%", "처리 중 iowait", BAD),
    ("61~88%", "같은 시각 CPU 유휴", MUTED),
    ("0.14", "DAQ 자체가 쓰는 코어 (12개 중)", MUTED),
], top)
bullets(s, [
    ("그래서 병렬 작업을 늘려도 거의 소용이 없다. DAQ 가 쓰고 있는 같은 "
     "파일시스템에 경합만 더할 뿐이다.", ""),
], y + Inches(0.3), size=16)
y2 = table(s, [
    ["Merged / PRD 를 어디에 쓰는가", "서브런당", "60초 대비 여유"],
    ["/scratch  (NFS)", "41.0 초", "1.5배"],
    ["/Data_ssd  (로컬 NVMe)", "+27.7 초", "+2.1배"],
], y + Inches(0.95), [Inches(6.6), Inches(2.6), Inches(2.6)])
footer(s, "원시 데이터는 NFS 에 그대로 둔다. 산출물만 옮긴다. 매크로는 그대로 — 디렉터리를 심볼릭 링크로 걸었다.")

# =========================================================================
# 18  divider - operations
# =========================================================================
divider(4, "운용", "운영자가 실제로 보는 화면과 치는 명령.")

# =========================================================================
# 19  screen
# =========================================================================
s, top = new("운용", "화면 하나, pane 다섯")
codebox(s, [
    "+-------------------------+------------------+",
    "| DAQ Run Status(monitor) |                  |  왼쪽  DAQ 상태",
    "|                     28  |   work space  7  |  오른쪽 작업용 셸과",
    "+-------------------------+                  |         데이터 이동",
    "| supervisor           8  |                  |",
    "+-------------------------+------------------+  ssh 끊김, 화면 죽음,",
    "| postrun              5  | dataflow:     3  |  터미널 재시작에도",
    "|                         |  ssd->data->     |  살아남는다",
    "|                         |  khu->scratch    |",
    "+-------------------------+------------------+",
    "            46%                    54%",
], M, top, W - 2 * M, size=14)
bullets(s, [
    ("감시자는 rcterm 을 한 줄 모드로 강제한다. 전체 화면은 매번 화면을 지워서 "
     "감시자 메시지까지 덮어버리기 때문이다. monitor pane 은 그 화면을 "
     "heartbeat 파일로 다시 그린다 — 읽기 전용이라 몇 개를 띄워도 되고, "
     "heartbeat 나이를 함께 보여줘서 멈춘 런이 그냥 정지 화면이 아니라 "
     "멈춘 것으로 보인다.", ""),
], Inches(4.55), size=15.5)

# =========================================================================
# 20  commands
# =========================================================================
s, top = new("운용", "실제로 쓰는 명령")
codebox(s, [
    "# 화면 전체 띄우기",
    "scripts/daq-tmux.sh                 # 화면만. 하드웨어 미접촉",
    "scripts/daq-tmux.sh --start         # DAQ 가 꺼져 있으면 같이 기동",
    "",
    "# 나중에 어디서든 다시 붙기",
    "tmux attach -t daq                  # Ctrl-B 다음 D 로 분리",
    "Ctrl-B  =                           # pane 비율 복원",
    "",
    "# DAQ 안전 종료 — 감시자 PID 에만 보낸다",
    "kill -TERM $(pgrep -x rcsupervisor)",
    "",
    "# 후처리",
    "scripts/postrun.sh --once --dry-run           # 무엇을 할지만 확인",
    "scripts/postrun.sh --follow --outroot /Data_ssd/RAW",
], M, top, W - 2 * M, size=14)
bullets(s, [
    ("프로세스 그룹에 신호를 보내지 말 것. daq, tcb, merger 는 rcterm 의 자식이 "
     "아니다. 그룹에 보내면 쓰기 도중에 끊겨 마지막 파일이 상한다.", ""),
], Inches(5.75), size=15)

# =========================================================================
# 21  storage
# =========================================================================
s, top = new("운용", "하루 운용의 비용")
y = table(s, [
    ["서브런 1개 (60초 분량)", "크기"],
    ["원시  FADC + SADC", "78 MB"],
    ["merged", "80 MB"],
    ["produced (PRD)", "77 MB"],
    ["합계", "235 MB"],
], top, [Inches(6.0), Inches(3.0)])
metric_row(s, [
    ("334 GB", "24시간 런 1회당", ACCENT2),
    ("48일", "/scratch 여유", MUTED),
    ("9일", "/Data_ssd 여유", BAD),
], y + Inches(0.4))
bullets(s, [
    ("로컬 디스크는 쓰기는 빠르지만 담아둘 자리는 좁다. 보관 정책은 차기 "
     "전에 정해야 한다.", ""),
], y + Inches(2.2), size=16)

# =========================================================================
# 22  status
# =========================================================================
s, top = new("현황", "지금 어디까지 왔나")
y = bullets(s, [
    ("수정본으로 24시간 로테이션을 두 번 통과했다. 둘 다 ENDED 가 찍히고 종료코드 0 "
     "으로 끝났으며 카탈로그가 완전히 마감됐다. 한 번은 우연일 수 있지만 두 번은 재현이다.", "good"),
    ("1.5일치를 전수 점검한 결과 고칠 결함이 없었다. health 286회 전부 통과, "
     "복구·FATAL 0건, 좀비 서브런 0건, 고아 행 0건.", "good"),
    ("두 런 모두 FADC·SADC·Merged·PRD 가 1440 으로 정확히 일치 — 어느 단계에서도 "
     "누락이 없었다.", "good"),
    ("후처리가 수집을 3개 서브런 뒤에서 그대로 따라붙고 있다.", "good"),
], top, size=16)
bullets(s, [
    ("남은 것은 정확성이 아니라 용량이다", "head"),
    ("런 하나가 로컬 디스크에 산출물 약 217 GB 를 남기는데 여유가 1.5 TB 다. 약 일주일치다. "
     "끝난 런은 NFS 로 되돌려야 하고, 심볼릭 링크 때문에 손으로 하기 번거로우니 "
     "드라이버에 정리 옵션을 붙이는 편이 낫다.", ""),
], y + Inches(0.2), size=16)

# =========================================================================
# 23  close
# =========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, PANEL)
rect(s, 0, 0, Inches(0.28), H, ACCENT)
tf = _tb(s, Inches(1.2), Inches(1.5), Inches(10.8), Inches(4.5))
_p(tf, "한 줄로 요약하면", 13, ACCENT, bold=True, font=MONO, first=True)
_p(tf, "이제 데이터 수집이 데스크톱 세션에 매이지 않고,",
   28, INK, bold=True, space_before=14, space_after=2)
_p(tf, "로테이션이 런 기록을 잃지 않으며,", 28, INK, bold=True, space_after=2)
_p(tf, "런이 끝나면 분석 파일이 준비돼 있다.", 28, INK, bold=True, space_after=22)
_p(tf, "github.com/Sang-Yong/RENE-daq-rcterm", 16, MUTED, font=MONO)
_p(tf, "docs/MANUAL.md   ·   docs/POSTRUN.md   ·   config/dotfiles/README.md",
   16, MUTED, font=MONO, space_before=4)

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "deck.pptx"
prs.save(out)
print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {out}")
